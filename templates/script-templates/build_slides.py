#!/usr/bin/env python3
"""
build_slides.py - Generates presentation slide decks (Marp Markdown & Standalone Web Presentation HTML)
for Workshop Harness projects.
"""

import sys
import os
import re
import argparse
import shutil
import subprocess
from pathlib import Path

def generate_slides(target_dir: str, output_dir: str = None, export_pdf: bool = False):
    target = Path(target_dir).resolve()
    if not target.exists():
        print(f"❌ Error: Workshop directory '{target}' does not exist.")
        return False

    out_path = Path(output_dir).resolve() if output_dir else target / "output" / "slides"
    out_path.mkdir(parents=True, exist_ok=True)

    # Read workshop title and topic
    readme_path = target / "README.md"
    title = target.name
    topic = "Hands-on Technical Workshop"

    if readme_path.exists():
        content = readme_path.read_text(encoding="utf-8")
        title_match = re.search(r"^#\s+(.+)$", content, re.MULTILINE)
        if title_match:
            title = title_match.group(1).strip()
        topic_match = re.search(r"> Topic:\s*(.+)$", content, re.MULTILINE)
        if topic_match:
            topic = topic_match.group(1).strip()

    # Read lab sections
    labs_file = target / "workshop" / "03_labs" / "README.md"
    lab_sections = []
    if labs_file.exists():
        labs_content = labs_file.read_text(encoding="utf-8")
        raw_sections = re.split(r"(?=\n##\s+)", "\n" + labs_content)
        for sec in raw_sections:
            sec = sec.strip()
            if sec:
                lab_sections.append(sec)

    if not lab_sections:
        lab_sections = [
            "## Lab 01: Environment & Hello World\nVerify local AI environment and run initial sanity check.",
            "## Lab 02: Core Hands-on Logic\nImplement structured outputs, schema verification, and core pipeline.",
            "## Lab 03: Full Pipeline & Challenge\nRun end-to-end evaluation, benchmark metrics, and explore bonus challenge."
        ]

    print(f"🎨 Generating Presentation Slide Decks for '{title}'...")

    # 1. Generate Marp Markdown (slides.md)
    marp_content = f"""---
marp: true
theme: default
paginate: true
header: "**{title}** | Hands-on Workshop"
footer: "Build with AI / DevFest Technical Workshop"
style: |
  section {{
    font-family: 'Inter', -apple-system, sans-serif;
    font-size: 26px;
    padding: 40px 60px;
    background-color: #0f172a;
    color: #f8fafc;
  }}
  h1 {{
    color: #38bdf8;
    font-size: 44px;
  }}
  h2 {{
    color: #818cf8;
    font-size: 34px;
    border-bottom: 2px solid #334155;
    padding-bottom: 10px;
  }}
  code {{
    background: #1e293b;
    color: #38bdf8;
  }}
  pre {{
    background: #1e293b;
    border: 1px solid #334155;
    border-radius: 8px;
  }}
  .badge {{
    background: #4f46e5;
    color: white;
    padding: 4px 12px;
    border-radius: 9999px;
    font-size: 16px;
  }}
---

<!-- _class: lead -->
<!-- _paginate: false -->
# {title}

### {topic}

<span class="badge">Hands-on Technical Lab</span>

**Facilitator**: Workshop Team  
**Event**: Build with AI / Community Coding Session

---

## 📋 Session Agenda & Timeline

* **00:00 - 00:10** (10m) ➔ Welcome & Architecture Overview
* **00:10 - 00:20** (10m) ➔ Environment & Hardware Verification
* **00:20 - 00:45** (25m) ➔ **Lab 01 & Lab 02 Hands-on**
* **00:45 - 01:10** (25m) ➔ **Lab 03 Pipeline & Challenge**
* **01:10 - 01:20** (10m) ➔ Wrap-up, Q&A & Community Resources

<!-- Speaker Notes: Greet attendees, check that everyone is seated and on the workshop Wi-Fi. -->

---

## ⚡ Prerequisites & Environment Check

* **Local Environment Setup**:
  * Run: `./scripts/check_env.sh` (Windows: `.\\scripts\\check_env.ps1`)
* **Google Colab Fallback (Zero Setup)**:
  * For low-spec laptops or GPU constraints, open the Colab notebook directly:
  * `output/colab/workshop_starter.ipynb`
* **API Keys & Credentials**:
  * Set `GEMINI_API_KEY` or local model server (`localhost:11434`)

<!-- Speaker Notes: Instruct TAs to roam and help attendees whose check_env script reports missing dependencies. -->

---

## 🏗️ Architecture & Core Concepts

```text
+-----------------------+      +-----------------------+      +-----------------------+
|   Attendee Input      | ---> |  Local LLM / Gemini   | ---> |  Structured Result    |
|  (Natural Language)   |      |  (Gemma 4 / Ollama)   |      |  (Pydantic / JSON)    |
+-----------------------+      +-----------------------+      +-----------------------+
                                           |
                                           v
                               +-----------------------+
                               |   Evaluation & Test   |
                               +-----------------------+
```

---
"""

    for idx, lab in enumerate(lab_sections, 1):
        lab_title_match = re.match(r"^##\s+(.+)$", lab, re.MULTILINE)
        lab_title = lab_title_match.group(1).strip() if lab_title_match else f"Lab {idx:02d}"
        clean_body = re.sub(r"^##\s+.+$", "", lab, flags=re.MULTILINE).strip()

        marp_content += f"""
## 🚀 {lab_title}

{clean_body[:400]}...

```bash
# Run and verify your Lab {idx:02d} implementation:
python3 workshop/01_starter/main.py
```

* 💡 **Tip**: Refer to `workshop/02_final/main.py` if you get stuck!
* ☁️ **Colab Badge**: Open in Google Colab for instant GPU execution.

---
"""

    marp_content += f"""
## 🚨 Live Troubleshooting & Emergency Fallbacks

| Symptom / Error | Quick 10s Hotfix Command |
| :--- | :--- |
| **Port 11434 Busy** | `pkill ollama && ollama serve &` |
| **OOM / Low RAM** | Fallback to Google Colab Starter Notebook |
| **API Key Missing** | `export GEMINI_API_KEY="YOUR_KEY"` |

👉 Ask roaming TAs for immediate live assistance!

---

<!-- _class: lead -->
# 🎉 Congratulations & Wrap-up!

### You've successfully built and verified the workshop pipeline!

* 📂 **GitHub Repo**: [https://github.com/JAICHANGPARK/{target.name}](https://github.com/JAICHANGPARK/{target.name})
* 📄 **PDF Handout**: `output/pdf/`
* ☁️ **Colab Notebooks**: `output/colab/`
* 💬 **Feedback & Q&A**: Please take 1 minute to fill out the session survey!
"""

    marp_file = out_path / "slides.md"
    marp_file.write_text(marp_content, encoding="utf-8")

    # 2. Generate Standalone Web Presentation HTML (index.html)
    slides_data = [
        {
            "title": title,
            "subtitle": topic,
            "badge": "Hands-on Technical Lab",
            "content": f"<p style='margin-top:20px; font-size:1.2rem; color:#94a3b8;'>Welcome to the session! Press <strong>Right Arrow</strong> or <strong>Space</strong> to advance.</p>"
        },
        {
            "title": "📋 Session Agenda & Timeline",
            "subtitle": "Minute-by-minute session flow",
            "badge": "Schedule",
            "content": "<ul><li><strong>00:00 - 00:10</strong> (10m) ➔ Welcome & Architecture Overview</li><li><strong>00:10 - 00:20</strong> (10m) ➔ Environment & Hardware Check</li><li><strong>00:20 - 00:45</strong> (25m) ➔ <strong>Lab 01 & 02 Hands-on</strong></li><li><strong>00:45 - 01:10</strong> (25m) ➔ <strong>Lab 03 Pipeline & Challenge</strong></li><li><strong>01:10 - 01:20</strong> (10m) ➔ Wrap-up, Q&A & Survey</li></ul>"
        },
        {
            "title": "⚡ Prerequisites & Environment Check",
            "subtitle": "Ensure zero-setup local & cloud readiness",
            "badge": "Setup",
            "content": "<ul><li>Run verification script: <code>./scripts/check_env.sh</code></li><li>Google Colab Cloud GPU Fallback: <code>output/colab/workshop_starter.ipynb</code></li><li>Set credentials: <code>export GEMINI_API_KEY=...</code></li></ul>"
        }
    ]

    for idx, lab in enumerate(lab_sections, 1):
        lab_title_match = re.match(r"^##\s+(.+)$", lab, re.MULTILINE)
        lab_title = lab_title_match.group(1).strip() if lab_title_match else f"Lab {idx:02d}"
        clean_body = re.sub(r"^##\s+.+$", "", lab, flags=re.MULTILINE).strip()
        slides_data.append({
            "title": f"🚀 {lab_title}",
            "subtitle": f"Hands-on Lab Step {idx}",
            "badge": f"Lab {idx:02d}",
            "content": f"<p>{clean_body[:300]}...</p><pre><code># Test command:\npython3 workshop/01_starter/main.py</code></pre>"
        })

    slides_data.append({
        "title": "🚨 Live Troubleshooting & Hotfixes",
        "subtitle": "Quick resolution for common errors",
        "badge": "Troubleshooting",
        "content": "<table style='width:100%; border-collapse:collapse; margin-top:15px;'><tr><th style='text-align:left; border-bottom:1px solid #334155; padding:8px;'>Issue</th><th style='text-align:left; border-bottom:1px solid #334155; padding:8px;'>Quick Hotfix</th></tr><tr><td style='padding:8px;'>Port 11434 Busy</td><td style='padding:8px;'><code>pkill ollama && ollama serve &</code></td></tr><tr><td style='padding:8px;'>Low RAM / Crash</td><td style='padding:8px;'>Switch to Google Colab Notebook</td></tr><tr><td style='padding:8px;'>API Key Missing</td><td style='padding:8px;'><code>export GEMINI_API_KEY=...</code></td></tr></table>"
    })

    slides_data.append({
        "title": "🎉 Congratulations & Wrap-up!",
        "subtitle": "You've completed all workshop milestones",
        "badge": "Complete",
        "content": "<p style='font-size:1.2rem; line-height:1.8;'>Thank you for participating! Check out the generated PDF handout and Colab interactive notebooks in your repo.</p>"
    })

    html_slides = ""
    for idx, s in enumerate(slides_data, 1):
        active_class = "active" if idx == 1 else ""
        html_slides += f"""
        <div class="slide {active_class}" id="slide-{idx}" data-index="{idx}">
          <div class="slide-header">
            <span class="slide-badge">{s['badge']}</span>
            <span class="slide-number">{idx} / {len(slides_data)}</span>
          </div>
          <h1 class="slide-title">{s['title']}</h1>
          <h3 class="slide-subtitle">{s['subtitle']}</h3>
          <div class="slide-body">
            {s['content']}
          </div>
        </div>
        """

    html_content = f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>{title} - Presentation Slides</title>
  <link rel="preconnect" href="https://fonts.googleapis.com">
  <link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>
  <link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;600;700;800&family=JetBrains+Mono:wght@400;600&display=swap" rel="stylesheet">
  <style>
    :root {{
      --bg: #0b0f19;
      --card-bg: #111827;
      --text: #f8fafc;
      --text-muted: #94a3b8;
      --primary: #38bdf8;
      --accent: #818cf8;
      --border: #1e293b;
    }}
    * {{ box-sizing: border-box; margin: 0; padding: 0; }}
    body {{
      font-family: 'Inter', sans-serif;
      background-color: var(--bg);
      color: var(--text);
      display: flex;
      flex-direction: column;
      height: 100vh;
      overflow: hidden;
    }}
    .presentation-container {{
      flex: 1;
      display: flex;
      align-items: center;
      justify-content: center;
      padding: 40px;
      position: relative;
    }}
    .slide {{
      display: none;
      width: 100%;
      max-width: 1100px;
      min-height: 600px;
      background: var(--card-bg);
      border: 1px solid var(--border);
      border-radius: 20px;
      padding: 60px 80px;
      box-shadow: 0 25px 50px -12px rgba(0, 0, 0, 0.5);
      flex-direction: column;
      justify-content: space-between;
      animation: fadeIn 0.3s cubic-bezier(0.16, 1, 0.3, 1);
    }}
    .slide.active {{
      display: flex;
    }}
    @keyframes fadeIn {{
      from {{ opacity: 0; transform: translateY(10px) scale(0.98); }}
      to {{ opacity: 1; transform: translateY(0) scale(1); }}
    }}
    .slide-header {{
      display: flex;
      justify-content: space-between;
      align-items: center;
      margin-bottom: 20px;
    }}
    .slide-badge {{
      background: rgba(56, 189, 248, 0.15);
      color: var(--primary);
      border: 1px solid rgba(56, 189, 248, 0.3);
      padding: 4px 14px;
      border-radius: 9999px;
      font-size: 0.85rem;
      font-weight: 600;
      text-transform: uppercase;
      letter-spacing: 0.05em;
    }}
    .slide-number {{
      color: var(--text-muted);
      font-family: 'JetBrains Mono', monospace;
      font-size: 0.9rem;
    }}
    .slide-title {{
      font-size: 2.8rem;
      font-weight: 800;
      color: #fff;
      line-height: 1.2;
      margin-bottom: 8px;
    }}
    .slide-subtitle {{
      font-size: 1.4rem;
      font-weight: 600;
      color: var(--accent);
      margin-bottom: 30px;
    }}
    .slide-body {{
      flex: 1;
      font-size: 1.15rem;
      line-height: 1.7;
      color: var(--text-muted);
    }}
    .slide-body ul {{
      margin-left: 24px;
      margin-top: 10px;
    }}
    .slide-body li {{
      margin-bottom: 12px;
    }}
    code {{
      font-family: 'JetBrains Mono', monospace;
      background: #1e293b;
      color: var(--primary);
      padding: 2px 8px;
      border-radius: 6px;
      font-size: 0.95em;
    }}
    pre {{
      background: #0f172a;
      border: 1px solid var(--border);
      border-radius: 10px;
      padding: 16px;
      margin-top: 16px;
      overflow-x: auto;
    }}
    pre code {{
      background: none;
      padding: 0;
    }}
    .controls-bar {{
      height: 70px;
      background: #070a11;
      border-top: 1px solid var(--border);
      display: flex;
      align-items: center;
      justify-content: space-between;
      padding: 0 40px;
    }}
    .controls-left, .controls-right {{
      display: flex;
      align-items: center;
      gap: 12px;
    }}
    .btn-ctrl {{
      background: #1e293b;
      border: 1px solid #334155;
      color: #fff;
      padding: 8px 16px;
      border-radius: 8px;
      font-size: 0.9rem;
      font-weight: 600;
      cursor: pointer;
      transition: all 0.15s ease;
    }}
    .btn-ctrl:hover {{
      background: #334155;
      border-color: #475569;
    }}
    .progress-track {{
      flex: 1;
      max-width: 400px;
      height: 6px;
      background: #1e293b;
      border-radius: 9999px;
      overflow: hidden;
      margin: 0 20px;
    }}
    .progress-fill {{
      height: 100%;
      background: linear-gradient(90deg, var(--primary), var(--accent));
      width: 10%;
      transition: width 0.2s ease;
    }}
  </style>
</head>
<body>
  <div class="presentation-container">
    {html_slides}
  </div>

  <div class="controls-bar">
    <div class="controls-left">
      <span style="font-weight:700; color:var(--primary); font-size:0.95rem;">{title}</span>
    </div>

    <div class="progress-track">
      <div class="progress-fill" id="progressFill"></div>
    </div>

    <div class="controls-right">
      <button class="btn-ctrl" onclick="prevSlide()">◀ Prev (Left)</button>
      <button class="btn-ctrl" onclick="nextSlide()">Next (Right) ▶</button>
      <button class="btn-ctrl" onclick="toggleFullScreen()">⛶ Fullscreen (F)</button>
    </div>
  </div>

  <script>
    let currentSlide = 1;
    const totalSlides = {len(slides_data)};

    function showSlide(index) {{
      if (index < 1) index = 1;
      if (index > totalSlides) index = totalSlides;
      currentSlide = index;

      document.querySelectorAll('.slide').forEach(s => s.classList.remove('active'));
      const activeSlide = document.getElementById('slide-' + currentSlide);
      if (activeSlide) activeSlide.classList.add('active');

      const pct = (currentSlide / totalSlides) * 100;
      document.getElementById('progressFill').style.width = pct + '%';
    }}

    function nextSlide() {{
      if (currentSlide < totalSlides) showSlide(currentSlide + 1);
    }}

    function prevSlide() {{
      if (currentSlide > 1) showSlide(currentSlide - 1);
    }}

    function toggleFullScreen() {{
      if (!document.fullscreenElement) {{
        document.documentElement.requestFullscreen();
      }} else {{
        if (document.exitFullscreen) {{
          document.exitFullscreen();
        }}
      }}
    }}

    document.addEventListener('keydown', (e) => {{
      if (e.key === 'ArrowRight' || e.key === ' ' || e.key === 'PageDown') {{
        e.preventDefault();
        nextSlide();
      }} else if (e.key === 'ArrowLeft' || e.key === 'PageUp') {{
        e.preventDefault();
        prevSlide();
      }} else if (e.key === 'f' || e.key === 'F') {{
        toggleFullScreen();
      }} else if (e.key === 'Home') {{
        showSlide(1);
      }} else if (e.key === 'End') {{
        showSlide(totalSlides);
      }}
    }});

    showSlide(1);
  </script>
</body>
</html>
"""
    html_file = out_path / "index.html"
    html_file.write_text(html_content, encoding="utf-8")

    # 3. Generate Native 16:9 PPTX for Google Slides
    pptx_file = out_path / "slides.pptx"
    has_pptx = False
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE

        prs = Presentation()
        # Set 16:9 widescreen dimensions (Google Slides native)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_slide_layout = prs.slide_layouts[6]

        BG_COLOR = RGBColor(15, 23, 42)     # #0F172A Dark Slate
        TITLE_COLOR = RGBColor(56, 189, 248) # #38BDF8 Sky Cyan
        SUB_COLOR = RGBColor(129, 140, 248)  # #818CF8 Indigo
        TEXT_COLOR = RGBColor(241, 245, 249) # #F1F5F9 Slate White
        CARD_BG = RGBColor(30, 41, 59)      # #1E293B Card Slate

        def add_bg(slide):
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
            shape.fill.solid()
            shape.fill.fore_color.rgb = BG_COLOR
            shape.line.fill.background()

        # Slide 1: Title
        s1 = prs.slides.add_slide(blank_slide_layout)
        add_bg(s1)
        tb = s1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(40)
        p1.font.bold = True
        p1.font.color.rgb = TITLE_COLOR
        p1.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = f"\n{topic}"
        p2.font.size = Pt(22)
        p2.font.color.rgb = SUB_COLOR
        p2.alignment = PP_ALIGN.CENTER

        p3 = tf.add_paragraph()
        p3.text = "\nBuild with AI / DevFest Technical Workshop"
        p3.font.size = Pt(14)
        p3.font.color.rgb = TEXT_COLOR
        p3.alignment = PP_ALIGN.CENTER

        # Slide 2: Agenda
        s2 = prs.slides.add_slide(blank_slide_layout)
        add_bg(s2)
        tb = s2.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(5.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Workshop Agenda"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR

        agenda_items = [
            ("00 - 15m", "Architecture & Prerequisites Setup"),
            ("15 - 45m", "Lab 01 & Lab 02 Core Hands-on Labs"),
            ("45 - 55m", "Lab 03 Full Pipeline & Benchmark"),
            ("55 - 60m", "Q&A, Troubleshooting & Next Steps")
        ]
        for t_slot, t_name in agenda_items:
            p = tf.add_paragraph()
            p.text = f"• [{t_slot}] {t_name}"
            p.font.size = Pt(20)
            p.font.color.rgb = TEXT_COLOR

        # Lab Slides
        for i, sec in enumerate(lab_sections, 1):
            s = prs.slides.add_slide(blank_slide_layout)
            add_bg(s)
            tb = s.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(5.8))
            tf = tb.text_frame
            tf.word_wrap = True
            lines = [l.strip() for l in sec.split("\n") if l.strip()]
            sec_title = lines[0].replace("#", "").strip() if lines else f"Lab {i}"
            p = tf.paragraphs[0]
            p.text = sec_title
            p.font.size = Pt(30)
            p.font.bold = True
            p.font.color.rgb = TITLE_COLOR

            for line in lines[1:8]:
                p = tf.add_paragraph()
                p.text = f"• {line}"
                p.font.size = Pt(18)
                p.font.color.rgb = TEXT_COLOR

        # Final Wrap-up Slide
        s_end = prs.slides.add_slide(blank_slide_layout)
        add_bg(s_end)
        tb = s_end.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Congratulations! Lab Completed!"
        p.font.size = Pt(38)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = f"\nYou have successfully built and verified '{title}'."
        p2.font.size = Pt(20)
        p2.font.color.rgb = TEXT_COLOR
        p2.alignment = PP_ALIGN.CENTER

        prs.save(str(pptx_file))
        has_pptx = True
    except ImportError:
        pass

    # 4. Generate Google Apps Script (create_google_slides.gs)
    gas_content = f"""/**
 * Google Apps Script to build Google Slides directly in Google Drive.
 * 
 * Instructions:
 * 1. Open https://script.google.com/
 * 2. Create a New Project, paste this code, and click 'Run' -> 'createWorkshopSlides'.
 * 3. A new Google Slides presentation will be created in your Google Drive!
 */
function createWorkshopSlides() {{
  const title = "{title}";
  const topic = "{topic}";
  const deck = SlidesApp.create(title + " - Presentation");
  
  // Slide 1: Title
  const slide1 = deck.getSlides()[0];
  slide1.getBackground().setSolidFill("#0f172a");
  const titleShape = slide1.insertTextBox(title, 50, 150, 620, 100);
  titleShape.getText().getTextStyle().setFontSize(28).setForegroundColor("#38bdf8").setBold(true);
  const subShape = slide1.insertTextBox(topic + "\\nBuild with AI / DevFest Workshop", 50, 260, 620, 60);
  subShape.getText().getTextStyle().setFontSize(16).setForegroundColor("#818cf8");

  // Slide 2: Agenda
  const slide2 = deck.appendSlide(SlidesApp.PredefinedLayout.BLANK);
  slide2.getBackground().setSolidFill("#0f172a");
  const agTitle = slide2.insertTextBox("Workshop Agenda", 50, 40, 620, 50);
  agTitle.getText().getTextStyle().setFontSize(24).setForegroundColor("#38bdf8").setBold(true);
  const agBody = slide2.insertTextBox("• [00 - 15m] Architecture & Prerequisites Setup\\n• [15 - 45m] Lab 01 & Lab 02 Hands-on Implementation\\n• [45 - 55m] Lab 03 End-to-End Pipeline & Evaluation\\n• [55 - 60m] Q&A, Troubleshooting & Wrap-up", 50, 110, 620, 250);
  agBody.getText().getTextStyle().setFontSize(16).setForegroundColor("#f8fafc");

  Logger.log("✅ Google Slides created successfully: " + deck.getUrl());
}}
"""
    gas_file = out_path / "create_google_slides.gs"
    gas_file.write_text(gas_content, encoding="utf-8")

    # 5. Generate Google Slides Integration Guide & README.md
    slides_readme = f"""# Presentation Slides for {title}

> Topic: {topic}

---

## 📊 1. Google Slides (구글 프레젠테이션) 사용 방법

### 방법 A: PPTX 1초 가져오기 (가장 추천)
1. **Google Slides ([slides.google.com](https://slides.google.com))** 접속
2. **`새 프레젠테이션 시작`** 클릭
3. 상단 메뉴에서 **`파일(File)` ➔ `슬라이드 가져오기(Import slides)`** 클릭
4. **`업로드(Upload)`** 탭에서 **`output/slides/slides.pptx`** 파일을 드래그하여 업로드
5. **`모든 슬라이드 선택(Select all slides)`** ➔ **`슬라이드 가져오기(Import slides)`** 클릭
6. 16:9 와이드스크린 다크 테마의 완성된 구글 슬라이드가 즉시 생성됩니다!

### 방법 B: Google Apps Script 매크로로 자동 생성
1. [script.google.com](https://script.google.com) 접속 후 **새 프로젝트** 생성
2. `output/slides/create_google_slides.gs` 파일의 코드를 복사하여 붙여넣기
3. 상단 메뉴의 **`실행(Run)`** 클릭
4. 구글 드라이브 내에 새 구글 프레젠테이션이 자동 생성됩니다.

---

## 🌐 2. 단독 웹 브라우저 프레젠테이션 (Zero-Setup)

- 별도 도구 설치 없이 브라우저에서 바로 발표할 수 있습니다:
```bash
open output/slides/index.html
```
- 단축키:
  - `Space` / `→` : 다음 슬라이드
  - `←` : 이전 슬라이드
  - `F` : 전체화면 (Fullscreen)
  - `Home` / `End` : 처음 / 마지막 슬라이드

---

## 📝 3. Marp Markdown & PDF 변환

- VS Code Marp 확장 프로그램을 통해 `slides.md`를 열어 발표 및 편집할 수 있습니다.
- Marp CLI를 통한 PDF 컴파일:
```bash
npx @marp-team/marp-cli@latest output/slides/slides.md --pdf -o output/slides/slides.pdf
```
"""
    (out_path / "README.md").write_text(slides_readme, encoding="utf-8")

    print(f"✨ Presentation Slides generated successfully at: {out_path}")
    print(f"  - {marp_file.name} (Marp Markdown)")
    print(f"  - {html_file.name} (Interactive Web Presentation)")
    if has_pptx:
        print(f"  - {pptx_file.name} (16:9 Google Slides / PowerPoint native .pptx)")
    print(f"  - {gas_file.name} (Google Apps Script for Slides)")
    print(f"  - README.md (with Google Slides import guide)")

    if export_pdf:
        if shutil.which("npx"):
            print("📄 Compiling slides to PDF via Marp CLI...")
            pdf_out = out_path / "slides.pdf"
            subprocess.run(["npx", "-y", "@marp-team/marp-cli@latest", str(marp_file), "--pdf", "-o", str(pdf_out)], check=False)
        else:
            print("💡 npx not found. Open index.html in browser and press Print to PDF.")
    return True

def main():
    parser = argparse.ArgumentParser(description="Generate Workshop Presentation Slides (Google Slides, Marp & HTML)")
    parser.add_argument("--target", required=True, help="Path to workshop project directory")
    parser.add_argument("--output", default=None, help="Target output directory")
    parser.add_argument("--export-pdf", action="store_true", help="Export to PDF using Marp CLI")

    args = parser.parse_args()
    generate_slides(args.target, args.output, args.export_pdf)

if __name__ == "__main__":
    main()
